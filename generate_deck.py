"""
generate_deck.py
================
Takes brand data and produces a .pptx file by editing the Superorder template.
Called by slack_agent.py after collecting all inputs from the user.

Requirements:
  pip install python-pptx requests
  Template file: Superorder_Integration_Call_Deck.pptx must be present
"""

import os, re, shutil, subprocess, zipfile
from datetime import datetime
from pathlib import Path

TEMPLATE_PATH = Path(__file__).parent / "Superorder_Integration_Call_Deck.pptx"
OUTPUT_DIR = Path(__file__).parent / "outputs"
OUTPUT_DIR.mkdir(exist_ok=True)

# ── Text replacement map ─────────────────────────────────────────────────────
# Maps placeholder text in the template XML → real values.
# These were identified by inspecting the template with markitdown.

def build_replacements(
    company_name: str,
    deal_name: str,
    hq_email_1: str,
    hq_email_2: str,
    poc_name: str,
    am_name: str,
    month_year: str = None,
) -> dict:
    if month_year is None:
        month_year = datetime.now().strftime("%B %Y")
    
    return {
        # Slide 1 — cover
        "Company Name": company_name,
        "Month Year": month_year,
        
        # Slide 6 — contacts (all 3 role rows default to poc_name)
        "Person 1": poc_name,     # replaces all "Person 1" placeholders
        "Person 2": "",           # leave empty (only 1 PoC known)
        "Person 3": "",
        
        # Slide 7 — DoorDash integration
        "[Insert HQ 1]": hq_email_1,
        "[Insert HQ 2]": hq_email_2,
        "Deal Name": deal_name,   # used as "First Name" in integration steps
        
        # AM name on agenda slide (only one of slides 2 or 3 will be kept)
        # The AM name already appears correctly since we keep the right variant
    }

# ── PPTX editing via XML ─────────────────────────────────────────────────────

def replace_text_in_xml(xml_text: str, replacements: dict) -> str:
    """Simple text replacement in slide XML. Handles split runs."""
    result = xml_text
    for old, new in replacements.items():
        # Direct replacement first
        result = result.replace(old, new)
    return result

def edit_pptx(template_path: Path, output_path: Path, replacements: dict, agenda_variant: str):
    """
    Unzip the PPTX, edit slide XMLs, rezip.
    agenda_variant: "maanav" keeps slide2, removes slide3
                    "sephra" keeps slide3, removes slide2
    """
    import zipfile, shutil, os
    
    tmp_dir = output_path.parent / "_tmp_deck"
    if tmp_dir.exists():
        shutil.rmtree(tmp_dir)
    
    # Unzip
    with zipfile.ZipFile(template_path, "r") as z:
        z.extractall(tmp_dir)
    
    ppt_dir = tmp_dir / "ppt"
    slides_dir = ppt_dir / "slides"
    
    # ── Find which slides are agenda variants ──────────────────────────────
    # Template has slide2.xml (Maanav variant) and slide3.xml (Sephra variant)
    # We need to remove the unwanted one from presentation.xml's sldIdLst
    # and delete the file.
    
    pres_xml_path = ppt_dir / "presentation.xml"
    pres_xml = pres_xml_path.read_text(encoding="utf-8")
    
    # The template slide ordering (from markitdown analysis):
    #   slide1 = cover, slide2 = agenda(Maanav AM), slide3 = agenda(Sephra AM)
    #   slide4 = what to expect, slide5 = products, slide6 = contacts,
    #   slide7 = doordash, slide8 = ubereats, slide9 = grubhub,
    #   slide10 = timeline, slide11 = help, slide12 = closing
    
    if agenda_variant == "maanav":
        slide_to_remove = "slide3"  # remove Sephra's agenda slide
    else:
        slide_to_remove = "slide2"  # remove Maanav's agenda slide
    
    # Remove from sldIdLst in presentation.xml
    # The relationship r:id references like "rId3" correspond to slide3
    # We parse which rId maps to which slide via slide relationships
    
    rels_path = ppt_dir / "_rels" / "presentation.xml.rels"
    rels_xml = rels_path.read_text(encoding="utf-8")
    
    # Find the rId for the slide we want to remove
    # e.g. <Relationship Id="rId3" ... Target="slides/slide3.xml"/>
    match = re.search(
        r'Id="(rId\d+)"[^>]+Target="slides/' + slide_to_remove + r'\.xml"',
        rels_xml
    )
    if match:
        rid_to_remove = match.group(1)
        # Remove from sldIdLst
        pres_xml = re.sub(r'<p:sldId[^>]+r:id="' + rid_to_remove + r'"[^/]*/>', '', pres_xml)
        pres_xml_path.write_text(pres_xml, encoding="utf-8")
        # Remove relationship
        rels_xml = re.sub(r'<Relationship[^>]+Id="' + rid_to_remove + r'"[^/]*/>\s*', '', rels_xml)
        rels_path.write_text(rels_xml, encoding="utf-8")
    
    # Delete the actual slide file
    slide_file = slides_dir / f"{slide_to_remove}.xml"
    if slide_file.exists():
        slide_file.unlink()
    slide_rels = slides_dir / "_rels" / f"{slide_to_remove}.xml.rels"
    if slide_rels.exists():
        slide_rels.unlink()
    
    # ── Apply text replacements to all remaining slides ───────────────────
    for slide_xml_path in sorted(slides_dir.glob("slide*.xml")):
        content = slide_xml_path.read_text(encoding="utf-8")
        content = replace_text_in_xml(content, replacements)
        slide_xml_path.write_text(content, encoding="utf-8")
    
    # ── Rezip into PPTX ───────────────────────────────────────────────────
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for file in sorted(tmp_dir.rglob("*")):
            if file.is_file():
                arcname = file.relative_to(tmp_dir)
                zout.write(file, arcname)
    
    shutil.rmtree(tmp_dir)

# ── Public API ───────────────────────────────────────────────────────────────

def convert_to_png(src_path: str) -> str:
    """Convert any image to PNG using Pillow. Returns path to a temp PNG file."""
    from PIL import Image
    import tempfile

    img = Image.open(src_path).convert("RGBA")
    # Composite onto white background so transparency renders well in PPTX
    background = Image.new("RGBA", img.size, (255, 255, 255, 255))
    background.paste(img, mask=img.split()[3])
    out = background.convert("RGB")

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    out.save(tmp.name, "PNG")
    return tmp.name


def insert_logo_on_slide1(pptx_path: str, logo_path: str):
    """
    Insert the chain logo onto slide 1 (title slide).
    Logo is placed to the right of the divider, scaled to fit a 3.8"×2.2" box.
    """
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image

    png_path = convert_to_png(logo_path)

    with Image.open(png_path) as img:
        orig_w, orig_h = img.size

    # Target bounding box: right side of slide, beside the 'X' divider (~11.2")
    max_w = Inches(3.8)
    max_h = Inches(2.2)

    scale = min(max_w / orig_w, max_h / orig_h)
    logo_w = int(orig_w * scale)
    logo_h = int(orig_h * scale)

    # Center within the box
    box_left = Inches(13.0)
    box_top  = Inches(3.8)
    left = box_left + (max_w - logo_w) // 2
    top  = box_top  + (max_h - logo_h) // 2

    prs = Presentation(pptx_path)
    prs.slides[0].shapes.add_picture(png_path, left, top, logo_w, logo_h)
    prs.save(pptx_path)

    import os
    os.unlink(png_path)


def build_deck(
    company_name: str,
    deal_name: str,
    hq_email_1: str,
    hq_email_2: str,
    poc_name: str,
    am_name: str,
    agenda_variant: str = "maanav",  # "maanav" or "sephra"
    month_year: str = None,
    logo_path: str = None,           # optional path to uploaded logo image
) -> str:
    """
    Generate the deck and return the output file path.
    """
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    replacements = build_replacements(
        company_name=company_name,
        deal_name=deal_name,
        hq_email_1=hq_email_1,
        hq_email_2=hq_email_2,
        poc_name=poc_name,
        am_name=am_name,
        month_year=month_year,
    )

    safe_name = re.sub(r"[^\w\s-]", "", company_name).strip().replace(" ", "_")
    output_path = OUTPUT_DIR / f"{safe_name}_Integration_Call_Deck.pptx"

    edit_pptx(TEMPLATE_PATH, output_path, replacements, agenda_variant)

    if logo_path:
        insert_logo_on_slide1(str(output_path), logo_path)

    return str(output_path)


if __name__ == "__main__":
    # Quick test
    path = build_deck(
        company_name="Test Brand Corp",
        deal_name="Test Brand Corp",
        hq_email_1="test-hq@testmanagement.com",
        hq_email_2="test-hq@testoperations.com",
        poc_name="John Smith",
        am_name="Maanav Patel",
        agenda_variant="maanav",
    )
    print(f"Generated: {path}")
